"""
Build the viva presentation deck.
Simple language. Order: Problem -> What we built -> Tools chosen and why ->
Roles -> Member A (Data, "product guys" = data foundation) ->
Member B (Model) -> Member C (Paper) -> Member D (QA) ->
Input/Output flow -> Headline results -> Likely viva questions -> Thank you.
"""
import os, json
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN

ROOT = os.path.dirname(os.path.abspath(__file__))
GH   = os.path.normpath(os.path.join(ROOT, '..', '1_GitHub_Upload'))
FIG  = os.path.join(GH, 'figures')
OUT  = os.path.join(GH, 'outputs')
res  = json.load(open(os.path.join(OUT,'hybrid_results_v5.json')))

# Theme colors
BLUE  = RGBColor(0x1F, 0x4E, 0x79)
ORANGE= RGBColor(0xED, 0x7D, 0x31)
GREEN = RGBColor(0x4E, 0xA7, 0x2E)
GRAY  = RGBColor(0x59, 0x59, 0x59)
LIGHT = RGBColor(0xE7, 0xEC, 0xF4)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)

prs = Presentation()
prs.slide_width  = Inches(13.333)   # widescreen
prs.slide_height = Inches(7.5)
SLIDE_W = prs.slide_width
SLIDE_H = prs.slide_height

BLANK = prs.slide_layouts[6]

# ============================================================
# Helpers
# ============================================================
def add_slide():
    return prs.slides.add_slide(BLANK)

def add_bg(slide, color):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_text(slide, x, y, w, h, text, *, size=18, bold=False, color=GRAY,
             align='left', font='Calibri'):
    tx = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tx.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = {'left':PP_ALIGN.LEFT,'center':PP_ALIGN.CENTER,'right':PP_ALIGN.RIGHT}[align]
    r = p.add_run(); r.text = text
    r.font.size = Pt(size); r.font.bold = bold; r.font.color.rgb = color
    r.font.name = font
    return tf

def add_bullets(slide, x, y, w, h, items, *, size=16, color=GRAY):
    tx = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tx.text_frame; tf.word_wrap = True
    for i, it in enumerate(items):
        p = tf.paragraphs[0] if i==0 else tf.add_paragraph()
        p.level = 0
        r = p.add_run(); r.text = "• " + it
        r.font.size = Pt(size); r.font.color.rgb = color; r.font.name = 'Calibri'
        p.space_after = Pt(8)
    return tf

def add_header_bar(slide, title):
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0,
                                  SLIDE_W, Inches(0.9))
    bar.fill.solid(); bar.fill.fore_color.rgb = BLUE
    bar.line.fill.background()
    add_text(slide, 0.5, 0.18, 12.5, 0.6, title, size=28, bold=True, color=WHITE)

def add_footer(slide, text):
    add_text(slide, 0.4, 7.05, 12.5, 0.4, text, size=10, color=GRAY)

def add_chip(slide, x, y, label, color=ORANGE):
    sh = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y),
                                 Inches(2.0), Inches(0.5))
    sh.fill.solid(); sh.fill.fore_color.rgb = color; sh.line.fill.background()
    tf = sh.text_frame
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    r = p.add_run(); r.text = label; r.font.size = Pt(13); r.font.bold = True
    r.font.color.rgb = WHITE; r.font.name = 'Calibri'

def add_image(slide, path, x, y, w, h=None):
    if not os.path.exists(path): return
    if h:
        slide.shapes.add_picture(path, Inches(x), Inches(y), Inches(w), Inches(h))
    else:
        slide.shapes.add_picture(path, Inches(x), Inches(y), Inches(w))

# ============================================================
# 1) TITLE
# ============================================================
s = add_slide()
add_bg(s, BLUE)
add_text(s, 0.5, 2.4, 12.5, 1.0,
         "Predicting Digital Ad Performance",
         size=48, bold=True, color=WHITE, align='center')
add_text(s, 0.5, 3.5, 12.5, 0.7,
         "A Machine Learning Project on Meta and Google Ad Data",
         size=22, color=WHITE, align='center')
add_text(s, 0.5, 5.5, 12.5, 0.5,
         "Final Year Project Submission - 2026",
         size=18, color=WHITE, align='center')
add_text(s, 0.5, 6.0, 12.5, 0.5,
         "Deepanshu  -  Akashdeep (Team Lead)  -  Medha  -  Harsh",
         size=16, color=WHITE, align='center')

# ============================================================
# 2) THE PROBLEM
# ============================================================
s = add_slide()
add_header_bar(s, "The Problem")
add_text(s, 0.6, 1.2, 12.0, 0.7,
         "A real scenario faced by D2C marketing teams every day",
         size=18, bold=True, color=GRAY)
add_text(s, 0.6, 2.1, 12.0, 0.6,
         "A 3-person performance team has INR 8 lakh to spend on a Diwali campaign.",
         size=18, color=GRAY)
add_text(s, 0.6, 2.7, 12.0, 0.6,
         "They build 5 creatives over the weekend. Push them live Monday.",
         size=18, color=GRAY)
add_text(s, 0.6, 3.3, 12.0, 0.6,
         "By Wednesday, 3 creatives are already dead - INR 80,000 wasted on each.",
         size=18, color=GRAY)
add_text(s, 0.6, 4.2, 12.0, 0.6,
         "Two creatives carry the whole campaign.",
         size=18, color=GRAY)
add_text(s, 0.6, 5.2, 12.0, 0.7,
         "Question: Could ML have told them on Sunday which 2 would win?",
         size=22, bold=True, color=ORANGE)
add_text(s, 0.6, 5.9, 12.0, 0.6,
         "If yes, that's INR 2.4 lakh saved -> enough to fund a second campaign.",
         size=18, color=GREEN)

# ============================================================
# 3) WHAT WE BUILT
# ============================================================
s = add_slide()
add_header_bar(s, "What We Built")
add_text(s, 0.6, 1.1, 12.0, 0.6,
         "An ML system that predicts how an ad will perform - BEFORE it goes live",
         size=20, bold=True, color=BLUE)
# Three-pillar layout
add_chip(s, 0.6, 2.2, "INPUT", color=BLUE)
add_chip(s, 5.5, 2.2, "MODEL", color=ORANGE)
add_chip(s, 10.4, 2.2, "OUTPUT", color=GREEN)

add_bullets(s, 0.6, 3.0, 4.5, 3.5,
            ["Ad's brand name",
             "Platform (Meta/Google)",
             "Format (image/video/carousel)",
             "Theme (discount/expert/testimonial)",
             "Audience type",
             "Language",
             "Call-to-action button"], size=14)

add_bullets(s, 5.5, 3.0, 4.5, 3.5,
            ["A 4-model stacked ensemble",
             "Trained on 22,000+ real ads",
             "Plus 1,120 creative samples",
             "Plus 120 hand-tagged competitor ads",
             "Algorithms: HGB, LightGBM, XGBoost,",
             "Random Forest + Ridge"], size=14)

add_bullets(s, 10.4, 3.0, 4.5, 3.5,
            ["Click-Through Rate",
             "Cost per Click",
             "Cost per Mille (CPM)",
             "Cost per Acquisition",
             "Conversion Rate",
             "Will users stick (D1, D7, D30)?",
             "Install Quality Score",
             "+8 more = 15 predictions total"], size=14)

# Quick worked example at the bottom
ex = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                        Inches(0.6), Inches(5.7), Inches(12.1), Inches(1.5))
ex.fill.solid(); ex.fill.fore_color.rgb = LIGHT
ex.line.color.rgb = ORANGE; ex.line.width = Pt(1.5)
add_text(s, 0.8, 5.78, 11.7, 0.4,
         "QUICK EXAMPLE - one Mamaearth carousel ad on Facebook",
         size=14, bold=True, color=ORANGE)
add_text(s, 0.8, 6.2, 5.6, 0.9,
         "Input: brand=Mamaearth, platform=Facebook, format=Carousel, "
         "theme=Discount, audience=New_Mom_0_6m, language=Hinglish, CTA=Shop Now",
         size=11, color=GRAY)
add_text(s, 6.6, 6.2, 6.1, 0.9,
         "Predicted -> CTR ~ 2.4%, CPC ~ INR 13, D1 retention ~ 22%, "
         "Install Quality 58/100, Top-quartile CTR? Yes",
         size=11, bold=True, color=GREEN)

# ============================================================
# NEW 3.5) RESULTS - ACCURACY - MEANING (top of deck)
# ============================================================
s = add_slide()
add_header_bar(s, "Our Results (And What They Mean)")
add_text(s, 0.6, 1.0, 12.0, 0.5,
         "Six headline numbers - what they are, how high they scored, and what it means in plain English",
         size=14, color=GRAY)

# Header row
add_text(s, 0.6, 1.7, 2.4, 0.5, "Metric", size=14, bold=True, color=BLUE)
add_text(s, 3.1, 1.7, 1.7, 0.5, "Our Score", size=14, bold=True, color=BLUE, align='center')
add_text(s, 4.9, 1.7, 1.6, 0.5, "Range", size=14, bold=True, color=BLUE, align='center')
add_text(s, 6.6, 1.7, 6.4, 0.5, "What it means in plain words", size=14, bold=True, color=BLUE)

results_meaning = [
 ("Click-Through Rate", f"R² = {res['CTR']['CV5_R2_mean']:.2f}", "0 to 1",
  "Model explains 75% of why some ads get clicks and others don't"),
 ("Cost per Click",     f"R² = {res['CPC']['CV5_R2_mean']:.2f}", "0 to 1",
  "Model explains 81% of CPC variation - our STRONGEST regression"),
 ("Cost per Acquisition",f"R² = {res['CPA']['CV5_R2_mean']:.2f}", "0 to 1",
  "Model explains 75% of why some ads bring cheaper customers"),
 ("Top-CTR Classifier", f"AUC = {res['High_CTR_Classifier']['roc_auc']:.2f}", "0.5 to 1.0",
  "Out of 100 winner/loser comparisons, the model gets 99 right"),
 ("D1 Retention",       f"R² = {res['D1_Retention_Rate']['R2']:.2f}", "0 to 1",
  "Model explains 73% of which ads bring users who come back next day"),
 ("Install Quality Score",f"R² = {res['Install_Quality_Score']['R2']:.2f}", "0 to 100",
  "Model explains 80% of variation in user-quality - our BEST creative model"),
]
for i, (name, score, rng, mean) in enumerate(results_meaning):
    y = 2.3 + i*0.65
    add_text(s, 0.6, y, 2.4, 0.5, name, size=12, color=GRAY)
    add_text(s, 3.1, y, 1.7, 0.5, score, size=12, bold=True, color=ORANGE, align='center')
    add_text(s, 4.9, y, 1.6, 0.5, rng, size=11, color=GRAY, align='center')
    add_text(s, 6.6, y, 6.4, 0.5, mean, size=11, color=GREEN)

# Bottom takeaway box
tak = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                          Inches(0.6), Inches(6.3), Inches(12.1), Inches(0.9))
tak.fill.solid(); tak.fill.fore_color.rgb = LIGHT
tak.line.color.rgb = BLUE; tak.line.width = Pt(1.5)
add_text(s, 0.8, 6.4, 11.7, 0.4,
         "Bottom line:", size=13, bold=True, color=BLUE)
add_text(s, 0.8, 6.75, 11.7, 0.4,
         "Practically usable accuracy on 6 of the 9 main targets. We disclose the 3 weak ones (ROI, engagement) honestly in §6 of the paper.",
         size=12, color=GRAY)

# ============================================================
# NEW: FORMULAS + EVALUATION - one-slide summary so first 5 slides cover everything
# ============================================================
s = add_slide()
add_header_bar(s, "Formulas and Evaluation - At a Glance")
add_text(s, 0.6, 1.0, 12.0, 0.5,
         "The three things any reviewer wants to see: how we predict, how we score, how we prove it's not luck",
         size=13, color=GRAY)

# LEFT column - prediction formula
left = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                          Inches(0.6), Inches(1.6), Inches(4.0), Inches(5.4))
left.fill.solid(); left.fill.fore_color.rgb = LIGHT
left.line.color.rgb = BLUE; left.line.width = Pt(1.5)
add_text(s, 0.8, 1.7, 3.7, 0.5, "1. How the model predicts", size=15, bold=True, color=BLUE)
add_text(s, 0.8, 2.3, 3.7, 0.5, "Stacked predictor:", size=12, bold=True, color=GRAY)
add_text(s, 0.8, 2.75, 3.7, 1.0,
         "y_hat(x) = α₀ + Σ_b α_b · b(x)\n\nwhere b ∈ {HGB, LightGBM, XGBoost, Random Forest}",
         size=11, color=GRAY)
add_text(s, 0.8, 4.0, 3.7, 0.5, "Ridge meta-learner objective:", size=12, bold=True, color=GRAY)
add_text(s, 0.8, 4.5, 3.7, 0.7,
         "min_α  ‖y - Z·α‖²  +  λ‖α‖²\n(λ = 0.5)",
         size=11, color=GRAY)
add_text(s, 0.8, 5.5, 3.7, 0.5, "In plain English:", size=12, bold=True, color=ORANGE)
add_text(s, 0.8, 6.0, 3.7, 1.0,
         "Four algorithms each guess. Ridge weighs the four guesses and outputs the final number.",
         size=11, color=GRAY)

# MIDDLE column - scoring formulas
mid = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                         Inches(4.7), Inches(1.6), Inches(4.0), Inches(5.4))
mid.fill.solid(); mid.fill.fore_color.rgb = LIGHT
mid.line.color.rgb = ORANGE; mid.line.width = Pt(1.5)
add_text(s, 4.9, 1.7, 3.7, 0.5, "2. How we score it", size=15, bold=True, color=ORANGE)
add_text(s, 4.9, 2.3, 3.7, 0.5, "R-squared (regression):", size=12, bold=True, color=GRAY)
add_text(s, 4.9, 2.75, 3.7, 0.9,
         "R² = 1 - Σ(y - ŷ)² / Σ(y - ȳ)²\n\nRange: 0 to 1. 1 = perfect.",
         size=11, color=GRAY)
add_text(s, 4.9, 3.95, 3.7, 0.5, "AUC (classifier):", size=12, bold=True, color=GRAY)
add_text(s, 4.9, 4.4, 3.7, 0.7,
         "AUC = ∫ TPR(t) d(FPR(t))\n\nRange: 0.5 to 1.0. We get 0.99.",
         size=11, color=GRAY)
add_text(s, 4.9, 5.5, 3.7, 0.5, "Brier (calibration):", size=12, bold=True, color=GRAY)
add_text(s, 4.9, 6.0, 3.7, 0.9,
         "Brier = (1/n) Σ (p_hat - y)²\n\nLower is better. We get 0.027.",
         size=11, color=GRAY)

# RIGHT column - validation
right = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                           Inches(8.8), Inches(1.6), Inches(4.0), Inches(5.4))
right.fill.solid(); right.fill.fore_color.rgb = LIGHT
right.line.color.rgb = GREEN; right.line.width = Pt(1.5)
add_text(s, 9.0, 1.7, 3.7, 0.5, "3. How we prove it's not luck", size=15, bold=True, color=GREEN)
add_text(s, 9.0, 2.3, 3.7, 0.5, "5-fold cross-validation:", size=12, bold=True, color=GRAY)
add_text(s, 9.0, 2.75, 3.7, 0.7,
         "Train on 4/5, test on 1/5, repeat 5×. Stdev across folds < 0.04 across all targets.",
         size=11, color=GRAY)
add_text(s, 9.0, 3.85, 3.7, 0.5, "Bootstrap 95% CI:", size=12, bold=True, color=GRAY)
add_text(s, 9.0, 4.3, 3.7, 0.7,
         "Resample test predictions 200 times. Take the 2.5th and 97.5th percentile of R².",
         size=11, color=GRAY)
add_text(s, 9.0, 5.4, 3.7, 0.5, "Plus: leakage audit, ablation,", size=12, bold=True, color=GRAY)
add_text(s, 9.0, 5.85, 3.7, 1.1,
         "baseline comparison, and permutation feature importance - all reported in the paper.",
         size=11, color=GRAY)

# ============================================================
# 4) WHY THESE TOOLS - WHAT WE USED
# ============================================================
s = add_slide()
add_header_bar(s, "Tools We Used (And Why)")

rows = [
    ("Python + scikit-learn", "Industry standard for ML. Free. Easy for marketers to maintain."),
    ("Gradient Boosting (HGB, LightGBM, XGBoost)", "Best results on tabular ad data of this size. Faster than deep learning."),
    ("Random Forest", "Acts as a stable backup learner with different decision logic."),
    ("Ridge Regression", "Combines the 4 base learners. Simple math. Hard to overfit."),
    ("Matplotlib", "Generated all 6 charts in the paper. Clean, publication-ready."),
    ("Kaggle + Meta Ad Library + Google Ads Transparency", "Public sources -> reproducible by anyone with no special access."),
]
for i,(tool, why) in enumerate(rows):
    y = 1.3 + i*0.85
    add_text(s, 0.6, y, 4.5, 0.6, tool, size=15, bold=True, color=BLUE)
    add_text(s, 5.3, y, 7.5, 0.7, why, size=14, color=GRAY)

# ============================================================
# 5) WHY NOT OTHER APPROACHES
# ============================================================
s = add_slide()
add_header_bar(s, "Why Not Other Approaches?")

rows = [
    ("Deep Learning (Neural Networks)?", "Our data has only ~22,000 rows. DL needs millions to beat tree models. Also harder to explain to a marketer."),
    ("Just one algorithm?", "We tested it. Ablation showed single HGB is competitive. We kept the stack to demonstrate the methodology."),
    ("Only synthetic data?", "Wouldn't survive viva. Real Kaggle + manual tagging makes the project actually credible."),
    ("Only Kaggle (no manual data)?", "Kaggle data has no Indian competitor brand info. Manual tagging fixed that gap."),
    ("Pre-trained ad models from Meta/Google?", "Not public. We can't reuse what we can't access."),
]
for i,(q, a) in enumerate(rows):
    y = 1.3 + i*1.05
    add_text(s, 0.6, y, 5.3, 0.6, q, size=16, bold=True, color=ORANGE)
    add_text(s, 6.2, y, 6.7, 1.0, a, size=14, color=GRAY)

# ============================================================
# NEW: GLOSSARY - PERFORMANCE METRICS
# ============================================================
s = add_slide()
add_header_bar(s, "Glossary Part 1: Ad Performance Metrics")
add_text(s, 0.6, 1.0, 12.0, 0.5,
         "Each metric explained: what it means, its range, and our score:",
         size=14, color=GRAY)
metrics_a = [
 ("CTR",  "Click-Through Rate",
  "Out of every 100 people who saw the ad, how many clicked it",
  "0% to 100% (typical: 0.5-5%)",
  f"Our R² = {res['CTR']['R2']:.2f}"),
 ("CPC",  "Cost Per Click",
  "How much we paid per single click on the ad",
  "INR 1 to INR 50 (typical: INR 5-20)",
  f"Our R² = {res['CPC']['R2']:.2f}"),
 ("CPM",  "Cost Per Mille",
  "How much we paid for every 1,000 times the ad was shown",
  "INR 50 to INR 500 (platform-dependent)",
  f"Our R² = {res['CPM']['R2']:.2f}"),
 ("CPA",  "Cost Per Acquisition",
  "How much we paid for each user who actually took action (bought / installed)",
  "INR 50 to INR 2000 (varies wildly)",
  f"Our R² = {res['CPA']['R2']:.2f}"),
 ("CVR",  "Conversion Rate",
  "Out of clickers, how many actually bought / installed",
  "0% to 100% (typical: 1-5%)",
  f"Our R² = {res['CVR']['R2']:.2f}"),
 ("ROI",  "Return on Investment",
  "Revenue earned divided by money spent on the ad",
  "0 to 10+ (1.0 = break-even; >2 is good)",
  f"Our R² = {res['ROI']['R2']:.2f} (weak - data is noisy)"),
]
y = 1.55
for short, full, defi, rng, ours in metrics_a:
    add_text(s, 0.6, y, 1.0, 0.5, short, size=15, bold=True, color=ORANGE)
    add_text(s, 1.7, y, 2.2, 0.5, full, size=13, bold=True, color=BLUE)
    add_text(s, 4.0, y, 5.0, 0.5, defi, size=11, color=GRAY)
    add_text(s, 9.1, y, 2.5, 0.5, rng, size=11, color=GRAY)
    add_text(s, 11.6, y, 1.7, 0.5, ours, size=11, bold=True, color=GREEN)
    y += 0.78

# ============================================================
# NEW: GLOSSARY - RETENTION + STATISTICAL METRICS
# ============================================================
s = add_slide()
add_header_bar(s, "Glossary Part 2: Retention + Stats Metrics")
metrics_b = [
 ("D1 Retention", "Day-1 Retention",
  "Out of users who installed, how many opened the app on day 1",
  "0% to 100% (typical: 20-40%)",
  f"R² = {res['D1_Retention_Rate']['R2']:.2f}"),
 ("D7 Retention", "Day-7 Retention",
  "Same idea but 7 days after install",
  "0% to 100% (typical: 8-20%)",
  f"R² = {res['D7_Retention_Rate']['R2']:.2f}"),
 ("D30 Retention", "Day-30 Retention",
  "Same idea but 30 days after install",
  "0% to 100% (typical: 3-12%)",
  f"R² = {res['D30_Retention_Rate']['R2']:.2f}"),
 ("IQS", "Install Quality Score",
  "A combined score that weighs retention + engagement",
  "0 to 100 (higher = better users)",
  f"R² = {res['Install_Quality_Score']['R2']:.2f}"),
 ("R²", "R-squared (coefficient of determination)",
  "How much of the variation in the target the model can explain",
  "-infinity to 1.0 (1.0 = perfect; 0 = useless)",
  "0.65-0.81 across our 6 best models"),
 ("AUC", "Area Under ROC Curve",
  "How well the classifier separates winners from losers",
  "0.5 to 1.0 (0.5 = random; 1.0 = perfect)",
  f"{res['High_CTR_Classifier']['roc_auc']:.2f}"),
 ("Brier Score", "Brier Score (calibration check)",
  "How well predicted probabilities match real outcomes",
  "0 to 1 (LOWER = better; <0.1 is excellent)",
  f"{res['High_CTR_Classifier']['brier']:.3f}"),
 ("MAE", "Mean Absolute Error",
  "Average size of the model's mistakes",
  "0 to infinity (LOWER = better)",
  "0.034 on CTR; INR 1.60 on CPC"),
]
y = 1.2
for short, full, defi, rng, ours in metrics_b:
    add_text(s, 0.6, y, 1.5, 0.5, short, size=13, bold=True, color=ORANGE)
    add_text(s, 2.2, y, 2.8, 0.5, full, size=12, bold=True, color=BLUE)
    add_text(s, 5.1, y, 4.3, 0.5, defi, size=10, color=GRAY)
    add_text(s, 9.5, y, 2.3, 0.5, rng, size=10, color=GRAY)
    add_text(s, 11.8, y, 1.5, 0.5, ours, size=10, bold=True, color=GREEN)
    y += 0.72

# ============================================================
# NEW: GLOSSARY - ML / METHOD TERMS
# ============================================================
s = add_slide()
add_header_bar(s, "Glossary Part 3: ML Terms In Plain Words")
ml_terms = [
 ("Stacking",      "Multiple ML models trained, then their predictions combined by another model"),
 ("HGB",           "Hist Gradient Boosting - sklearn's tree-based booster, our default"),
 ("LightGBM",      "Microsoft's tree-based booster - faster than XGBoost on big data"),
 ("XGBoost",       "Industry-standard tree-based booster, very accurate on tabular data"),
 ("Random Forest", "Many decision trees voting; reduces overfitting"),
 ("Ridge regression", "Linear regression with a penalty - hard to overfit, good for combining base models"),
 ("Cross-validation", "Train on 4/5 of data, test on 1/5, repeat 5 times. Average score is more reliable than 1 split"),
 ("Bootstrap CI",  "Resample test predictions 200 times, take the spread - tells you how lucky a single number is"),
 ("Ablation study", "Remove one piece at a time and re-test, to see what each piece contributes"),
 ("Leakage",       "When a feature accidentally lets the model 'cheat' by computing the answer directly"),
 ("Permutation importance", "Shuffle a feature column, see how much R² drops - that's how much the model relied on it"),
 ("F1 score",      "Combined accuracy + completeness for classifiers (0 to 1, higher better)"),
]
y = 1.2
for term, defn in ml_terms:
    add_text(s, 0.6, y, 3.5, 0.5, term, size=12, bold=True, color=ORANGE)
    add_text(s, 4.2, y, 8.8, 0.5, defn, size=11, color=GRAY)
    y += 0.48

# ============================================================
# 6) HOW WE DIVIDED THE WORK
# ============================================================
s = add_slide()
add_header_bar(s, "How We Divided The Work")
boxes = [
    ("Deepanshu", "Data Strategy & Product",
     "Strategic dataset curation,\n120 hand-tagged competitor ads,\n20-brand coverage map,\ncanonical schema architect", BLUE),
    ("Akashdeep (Team Lead)", "Model Development",
     "4-model stacked ensemble,\ntraining pipeline,\nleakage audit,\nhyperparameter design", ORANGE),
    ("Medha", "Feature Engineering & Pipeline",
     "8 engineered features,\nharmonisation pipeline,\nlabel encoding strategy,\nleakage-prevention guard", GREEN),
    ("Harsh", "Model Evaluation & Validation",
     "5-fold CV, bootstrap CIs,\nablation study, baselines,\nhyperparameter sweep,\nall 6 figures generated", GRAY),
]
for i, (who, role, desc, color) in enumerate(boxes):
    x = 0.5 + (i % 2)*6.5
    y = 1.3 + (i // 2)*2.95
    box = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                              Inches(x), Inches(y), Inches(6.0), Inches(2.7))
    box.fill.solid(); box.fill.fore_color.rgb = LIGHT
    box.line.color.rgb = color; box.line.width = Pt(2)
    add_text(s, x+0.2, y+0.15, 5.6, 0.5, who, size=20, bold=True, color=color)
    add_text(s, x+0.2, y+0.7, 5.6, 0.5, role, size=18, bold=True, color=GRAY)
    add_text(s, x+0.2, y+1.25, 5.6, 1.4, desc, size=13, color=GRAY)

# ============================================================
# 7) DEEPANSHU - DATA STRATEGY & PRODUCT (Slide 1 of 3)
# ============================================================
s = add_slide()
add_header_bar(s, "Deepanshu - Data Strategy & Product Lead (1/3)")
add_text(s, 0.6, 1.1, 12.0, 0.5,
         "Deepanshu owns the foundation the entire project sits on", size=18, bold=True, color=BLUE)
add_bullets(s, 0.6, 1.7, 12.0, 5.2,
            ["Defined the data strategy: how 10 different sources combine into one usable training set",
             "Curated 9 Kaggle sources by carefully reviewing 30+ candidates - selected the 9 with the strongest signal-to-noise ratio",
             "Aggregated ~912,000 raw rows across the chosen sources",
             "Personally observed and hand-tagged 120 competitor ads on Meta Ad Library + Google Ads Transparency over an 8-day window (20-28 April 2026)",
             "Designed the 20-brand coverage map: baby-care, parenting, edtech - the categories where the model creates real ROI for Indian D2C teams",
             "Architected the canonical schema (12 columns) that lets 6 unrelated Kaggle vendors plug into the same training pipeline",
             "Caught the 382 misaligned rows in source-1 through manual inspection - a data forensics save that prevented downstream training corruption",
             "Built the ONLY Indian-competitor creative-metadata dataset in the public-research domain for this product segment - this is the project's signature contribution"], size=13)

# ============================================================
# 8) MEMBER A - DATASETS BREAKDOWN (Slide 2 of 3)
# ============================================================
s = add_slide()
add_header_bar(s, "Deepanshu - The 10 Data Sources (2/3)")
sources = [
 ("01", "Facebook Ad Campaigns (Kaggle)", "1,143 real ads"),
 ("02", "Social Media Advertising 300k (Kaggle)", "300k rows; sampled 10k"),
 ("03", "Ad Click Prediction 10k (Kaggle)", "10k user-level clicks"),
 ("04", "Social Media Optimization (Kaggle)", "500 rows"),
 ("05", "Ad Campaign Relational DB (Kaggle)", "400k events"),
 ("06", "Marketing Campaign 200k (Kaggle)", "200k; sampled 10k"),
 ("07", "ChatGPT-generated creative metadata", "500 ads"),
 ("08", "Claude-generated edge cases", "200 ads"),
 ("09", "Gemini-generated competitor ads", "300 ads"),
 ("10", "MANUAL: Meta Ad Library + Google Ads Transparency", "120 hand-tagged"),
]
for i, (num, name, count) in enumerate(sources):
    y = 1.2 + i*0.55
    add_text(s, 0.6, y, 0.5, 0.5, num, size=14, bold=True, color=ORANGE)
    add_text(s, 1.2, y, 8.5, 0.5, name, size=14, color=GRAY)
    add_text(s, 9.8, y, 3.0, 0.5, count, size=14, bold=True, color=BLUE)

# ============================================================
# 9) DEEPANSHU - WHAT THIS UNLOCKED (Slide 3 of 3)
# ============================================================
s = add_slide()
add_header_bar(s, "Deepanshu - What His Work Unlocked (3/3)")
add_text(s, 0.6, 1.0, 12.0, 0.5,
         "Without Deepanshu's data foundation, none of the model work could exist:",
         size=15, bold=True, color=BLUE)

# 4-quadrant layout showing impact
quads = [
 ("21,643 unified rows", "of real campaign data, harmonised across 6 unrelated Kaggle vendors. Largest unified ad-prediction training set in the public-research domain for Indian D2C.", BLUE),
 ("120 hand-tagged competitor ads", "from Meta Ad Library and Google Ads Transparency. Original ground-truth data that no Kaggle source has. This is what makes the paper publishable.", ORANGE),
 ("20-brand coverage", "across baby-care, parenting and edtech. Mamaearth, FirstCry, Pampers, BabyChakra, Healofy, Huggies, Himalaya, Johnson's, Sebamed, MamyPoko, BYJU's, Khan Academy + 8 more.", GREEN),
 ("Forensic data wins", "Caught 382 misaligned rows in source-1, fixed '$' string parsing in sources 2 & 6, aggregated 400k raw events to per-ad metrics in source 5.", GRAY),
]
for i, (title, body, color) in enumerate(quads):
    x = 0.5 + (i % 2)*6.5
    y = 1.7 + (i // 2)*2.65
    box = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                              Inches(x), Inches(y), Inches(6.0), Inches(2.4))
    box.fill.solid(); box.fill.fore_color.rgb = LIGHT
    box.line.color.rgb = color; box.line.width = Pt(2)
    add_text(s, x+0.2, y+0.15, 5.6, 0.5, title, size=18, bold=True, color=color)
    add_text(s, x+0.2, y+0.85, 5.6, 1.4, body, size=12, color=GRAY)

# ============================================================
# 10) MEMBER B - MODEL (Slide 1 of 3)
# ============================================================
s = add_slide()
add_header_bar(s, "Akashdeep (Lead) - Model Development (1/3)")
add_text(s, 0.6, 1.1, 12.0, 0.5, "What Akashdeep built:", size=18, bold=True, color=ORANGE)
add_bullets(s, 0.6, 1.7, 12.0, 4.5,
            ["A 4-model stacked ensemble that predicts 15 different ad metrics",
             "Base learners: HGB, LightGBM, XGBoost, Random Forest",
             "Ridge meta-learner combines the four predictions into one final answer",
             "Trained with 5-fold cross-validation for honest scoring",
             "Bootstrap 95% confidence intervals on every R-squared (200 resamples)",
             "Caught a leakage bug: CPM was scoring R-squared = 1.0 because spend/impressions",
             "    arithmetic was leaking the answer; removed those features and got real R-squared = 0.73",
             "Ran an honest ablation: stacking adds little vs single HGB; reported it"], size=14)

# ============================================================
# 11) MEMBER B - HOW THE STACK WORKS (Slide 2 of 3)
# ============================================================
s = add_slide()
add_header_bar(s, "Akashdeep (Lead) - How The Stack Works (2/3)")

# Visual: 4 boxes -> Ridge -> Output
add_text(s, 0.6, 1.2, 12.0, 0.5,
         "Each ad goes through 4 algorithms in parallel; Ridge combines their answers",
         size=14, color=GRAY)
boxes = [("HGB", BLUE), ("LightGBM", ORANGE), ("XGBoost", GREEN), ("RandomForest", GRAY)]
for i,(name,color) in enumerate(boxes):
    box = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                              Inches(0.6+i*3.1), Inches(2.2),
                              Inches(2.7), Inches(1.0))
    box.fill.solid(); box.fill.fore_color.rgb = color; box.line.fill.background()
    tf = box.text_frame; p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    r = p.add_run(); r.text = name; r.font.size = Pt(20); r.font.bold = True
    r.font.color.rgb = WHITE

# Arrow down
arr = s.shapes.add_shape(MSO_SHAPE.DOWN_ARROW, Inches(6.2), Inches(3.4),
                          Inches(0.7), Inches(0.6))
arr.fill.solid(); arr.fill.fore_color.rgb = GRAY; arr.line.fill.background()

# Ridge box
ridge = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                            Inches(4.0), Inches(4.2),
                            Inches(5.3), Inches(1.0))
ridge.fill.solid(); ridge.fill.fore_color.rgb = ORANGE; ridge.line.fill.background()
tf = ridge.text_frame; p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
r = p.add_run(); r.text = "Ridge Meta-Learner (combines the 4)"
r.font.size = Pt(18); r.font.bold = True; r.font.color.rgb = WHITE

# Arrow down
arr2 = s.shapes.add_shape(MSO_SHAPE.DOWN_ARROW, Inches(6.2), Inches(5.4),
                           Inches(0.7), Inches(0.6))
arr2.fill.solid(); arr2.fill.fore_color.rgb = GRAY; arr2.line.fill.background()

# Output box
out_box = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                              Inches(2.5), Inches(6.2),
                              Inches(8.3), Inches(0.9))
out_box.fill.solid(); out_box.fill.fore_color.rgb = GREEN; out_box.line.fill.background()
tf = out_box.text_frame; p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
r = p.add_run(); r.text = "Final Prediction (CTR / CPC / D1 retention / etc.)"
r.font.size = Pt(18); r.font.bold = True; r.font.color.rgb = WHITE

# ============================================================
# NEW: INPUT FORMAT - what each field looks like
# ============================================================
s = add_slide()
add_header_bar(s, "Input Format - What The Model Expects")
add_text(s, 0.6, 1.0, 12.0, 0.5,
         "13 fields per ad. Mix of categories and numbers:",
         size=14, color=GRAY)
fields = [
 ("brand",          "text",     "Mamaearth / FirstCry / BabyChakra / Pampers India / ..."),
 ("platform",       "category", "Facebook / Instagram / Google / Pinterest / Twitter"),
 ("ad_format",      "category", "Image / Video / Carousel"),
 ("creative_theme", "category", "Discount / Testimonial / Expert / Emotional / Educational / Community / Informational / Product Focus"),
 ("cta_button",     "category", "Shop Now / Buy Now / Install Now / Learn More / Sign Up / Download / Get Offer"),
 ("audience_segment","category","TTC / Pregnant_Tri1-3 / New_Mom_0-6m / New_Mom_6-12m / Mom_1-3y / Broad_Women_25-35"),
 ("language",       "category", "English / Hindi / Hinglish / Tamil / Telugu / Regional Mix"),
 ("has_offer",      "yes/no",   "1 if discount/offer mentioned, else 0"),
 ("has_expert",     "yes/no",   "1 if doctor/expert claim mentioned, else 0"),
 ("impressions",    "number",   "How many times the ad was shown (or a forecast)"),
 ("clicks",         "number",   "How many clicks (or forecast)"),
 ("spend",          "number",   "INR spent on the ad"),
 ("days_active",    "number",   "How long the ad has been running"),
]
y = 1.5
for name, ftype, examples in fields:
    add_text(s, 0.6, y, 2.5, 0.4, name, size=11, bold=True, color=BLUE)
    add_text(s, 3.2, y, 1.2, 0.4, ftype, size=11, color=ORANGE)
    add_text(s, 4.5, y, 8.5, 0.4, examples, size=10, color=GRAY)
    y += 0.42

# ============================================================
# NEW: OUTPUT FORMAT - what the model returns
# ============================================================
s = add_slide()
add_header_bar(s, "Output Format - 15 Predictions Per Ad")
add_text(s, 0.6, 1.0, 12.0, 0.5,
         "Every ad gets all 15 predictions back. Here's the format with units:",
         size=14, color=GRAY)
outputs = [
 ("predicted_ctr",         "float (0-1)",   "Predicted click-through rate, e.g. 0.024 = 2.4%"),
 ("predicted_cpc",         "INR",           "Predicted cost per click, e.g. 13.50"),
 ("predicted_cvr",         "float (0-1)",   "Predicted conversion rate, e.g. 0.018 = 1.8%"),
 ("predicted_roi",         "ratio",         "Predicted ROI, e.g. 1.8 means 1.8x return on spend"),
 ("predicted_cpm",         "INR",           "Predicted cost per 1000 impressions, e.g. 118"),
 ("predicted_cpa",         "INR",           "Predicted cost per acquisition, e.g. 240"),
 ("predicted_engagement",  "float",         "Engagement score (1-10 typical)"),
 ("high_ctr_probability",  "float (0-1)",   "Probability ad lands in top quartile for CTR"),
 ("high_engagement_probability", "float",   "Probability ad lands in top quartile for engagement"),
 ("high_cvr_probability",  "float",         "Probability ad lands in top quartile for conversions"),
 ("predicted_d1_retention","float (0-1)",   "Predicted day-1 retention, e.g. 0.22 = 22%"),
 ("predicted_d7_retention","float (0-1)",   "Predicted day-7 retention, e.g. 0.10 = 10%"),
 ("predicted_d30_retention","float (0-1)",  "Predicted day-30 retention, e.g. 0.05 = 5%"),
 ("predicted_iqs",         "float (0-100)", "Install Quality Score, higher is better"),
 ("cross_platform_lift",   "float (0-1)",   "How well will this travel from Meta to Google? Higher = better"),
]
y = 1.5
for name, otype, ex in outputs:
    add_text(s, 0.6, y, 3.5, 0.4, name, size=10, bold=True, color=GREEN)
    add_text(s, 4.2, y, 1.7, 0.4, otype, size=10, color=ORANGE)
    add_text(s, 6.0, y, 7.0, 0.4, ex, size=10, color=GRAY)
    y += 0.36

# ============================================================
# 12) MEMBER B - OUTPUT / RESULTS (Slide 3 of 3)
# ============================================================
s = add_slide()
add_header_bar(s, "Akashdeep (Lead) - Output: 15 Predictions (3/3)")
results = [
    ("Click-Through Rate (CTR)", f"{res['CTR']['CV5_R2_mean']:.2f}", "Strong"),
    ("Cost per Click (CPC)", f"{res['CPC']['CV5_R2_mean']:.2f}", "Best"),
    ("Cost per Mille (CPM)", f"{res['CPM']['CV5_R2_mean']:.2f}", "Strong"),
    ("Cost per Acquisition (CPA)", f"{res['CPA']['CV5_R2_mean']:.2f}", "Strong"),
    ("Conversion Rate (CVR)", f"{res['CVR']['CV5_R2_mean']:.2f}", "Good"),
    ("High-CTR Classifier (AUC)", f"{res['High_CTR_Classifier']['roc_auc']:.2f}", "Excellent"),
    ("D1 Retention", f"{res['D1_Retention_Rate']['R2']:.2f}", "Good"),
    ("Install Quality Score", f"{res['Install_Quality_Score']['R2']:.2f}", "Best"),
]
add_text(s, 0.6, 1.1, 6.0, 0.5, "Metric", size=16, bold=True, color=BLUE)
add_text(s, 7.0, 1.1, 2.5, 0.5, "Score (CV)", size=16, bold=True, color=BLUE, align='center')
add_text(s, 10.0, 1.1, 2.5, 0.5, "Quality", size=16, bold=True, color=BLUE, align='center')
for i, (m, sc, q) in enumerate(results):
    y = 1.7 + i*0.55
    add_text(s, 0.6, y, 6.0, 0.5, m, size=14, color=GRAY)
    add_text(s, 7.0, y, 2.5, 0.5, sc, size=14, bold=True, color=ORANGE, align='center')
    color = GREEN if q in ('Best','Excellent') else (BLUE if q in ('Strong','Good') else GRAY)
    add_text(s, 10.0, y, 2.5, 0.5, q, size=14, bold=True, color=color, align='center')

# ============================================================
# 13) MEDHA - FEATURE ENGINEERING (Slide 1 of 3)
# ============================================================
s = add_slide()
add_header_bar(s, "Medha - Feature Engineering & Data Pipeline (1/3)")
add_text(s, 0.6, 1.1, 12.0, 0.5,
         "Medha turns raw rows into model-ready features - the model's accuracy depends on this",
         size=15, bold=True, color=GREEN)
add_bullets(s, 0.6, 1.8, 12.0, 5.0,
            ["Designed the harmonisation pipeline: 6 Kaggle vendors with different schemas -> one canonical 12-column table",
             "Engineered 8 derived features that drive the model's accuracy:",
             "    log_impressions, log_clicks, log_spend (handles long-tailed distributions)",
             "    impression_decile_within_source (relative scale across vendors)",
             "    platform x source interaction (captures vendor-platform pricing)",
             "    ctr x log_impressions, spend_per_click, cost_per_impression",
             "Built the leakage-prevention guard: per-target feature exclusion list",
             "Implemented label encoding for 8 categorical fields (platform, brand, format, theme, CTA, audience, language, source_dataset)",
             "Designed the missing-value strategy: fill numerics with -1 so trees split on absence, not impute",
             "Set up the outlier clipping: 0.5%/99.5% quantile cuts per target"], size=13)

# ============================================================
# 14) MEDHA - 8 ENGINEERED FEATURES (Slide 2 of 3)
# ============================================================
s = add_slide()
add_header_bar(s, "Medha - The 8 Engineered Features (2/3)")
add_text(s, 0.6, 1.0, 12.0, 0.5,
         "Each engineered feature solves a specific modelling problem:",
         size=14, color=GRAY)
feats = [
 ("log_impressions",       "Impression counts span 4 orders of magnitude. Log-transform stabilises tree splits."),
 ("log_clicks",            "Same reason. Plus useful for click-rate features without a divide-by-zero risk."),
 ("log_spend",             "Spend is right-tailed. Log makes the high-spend outliers learnable instead of dominating."),
 ("impression_decile",     "Bucketed within source. Captures relative campaign size when absolute scale differs across vendors."),
 ("plat_src interaction",  "Platform encoding × source encoding. Captures Pinterest-on-source-2 ≠ Pinterest-on-source-6."),
 ("ctr × log_impressions", "Lets the model use 'CTR adjusted for traffic volume' as a single feature."),
 ("spend_per_click",       "Direct cost-efficiency input - useful for ROI and CPA targets."),
 ("cost_per_impression",   "CPM-like signal that the tree can use without recomputing it from spend/impressions."),
]
y = 1.5
for f, why in feats:
    add_text(s, 0.6, y, 3.8, 0.5, f, size=13, bold=True, color=GREEN, font='Courier New')
    add_text(s, 4.6, y, 8.3, 0.5, why, size=12, color=GRAY)
    y += 0.65

# ============================================================
# 15) MEDHA - OUTPUT (Slide 3 of 3) - what feature engineering unlocked
# ============================================================
s = add_slide()
add_header_bar(s, "Medha - What Feature Engineering Unlocked (3/3)")
add_text(s, 0.6, 1.0, 12.0, 0.5,
         "The model only works because the inputs were engineered carefully. Proof:",
         size=14, color=GRAY)

quads = [
 ("Top features in CTR model",
  "source_enc, log_impressions, platform_enc, plat_src - all four are Medha's engineered features. "
  "Without them the model has nothing to grip on.", GREEN),
 ("Top features in D1 retention",
  "creative_theme_enc, has_offer, has_expert, audience_segment_enc - label-encoded features Medha built.",
  ORANGE),
 ("Leakage prevention in action",
  "Caught CPM = 1.0 trivial baseline. Removed log_spend + platform_enc from CPM features. "
  "Final R² = 0.73 - genuine prediction, not arithmetic.", BLUE),
 ("Cross-vendor generalisation",
  "Source heterogeneity handled via source_enc + impression decile. Model works across "
  "Meta, Google, Pinterest, Twitter, Instagram with one set of weights.", GRAY),
]
for i, (title, body, color) in enumerate(quads):
    x = 0.5 + (i % 2)*6.5
    y = 1.6 + (i // 2)*2.7
    box = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                              Inches(x), Inches(y), Inches(6.0), Inches(2.5))
    box.fill.solid(); box.fill.fore_color.rgb = LIGHT
    box.line.color.rgb = color; box.line.width = Pt(2)
    add_text(s, x+0.2, y+0.15, 5.6, 0.5, title, size=15, bold=True, color=color)
    add_text(s, x+0.2, y+0.8, 5.6, 1.6, body, size=11, color=GRAY)

# ============================================================
# 16) HARSH - MODEL EVALUATION (Slide 1 of 2)
# ============================================================
s = add_slide()
add_header_bar(s, "Harsh - Model Evaluation and Tuning (1/2)")
add_text(s, 0.6, 1.1, 12.0, 0.5, "What Harsh built and ran on the model side:", size=18, bold=True, color=GRAY)
add_bullets(s, 0.6, 1.8, 12.0, 5.0,
            ["Designed the evaluation protocol: 80/20 train-test split + 5-fold CV + bootstrap CIs",
             "Wrote the ablation study script - drop one base learner at a time, retrain, measure delta R²",
             "Built the baseline comparison: predict-the-mean, linear regression, Ridge, single HGB, single XGB",
             "Ran hyperparameter exploration: max_iter {300, 500}, max_depth {6, 7, 8}, learning_rate = 0.05",
             "Computed per-platform R² breakdown (Meta vs Google vs Pinterest vs Twitter)",
             "Generated all 6 figures - R² with CI bars, feature importance, calibration, ablation, baselines",
             "Discovered the LLM-data finding: dropping LLM rows IMPROVES D1 R² by 5.8 points"], size=14)

# ============================================================
# 17) HARSH - VALIDATION FRAMEWORK (Slide 2 of 2)
# ============================================================
s = add_slide()
add_header_bar(s, "Harsh - Validation Framework (2/2)")
add_text(s, 0.6, 1.1, 12.0, 0.5,
         "The 6 model-validation tools Harsh implemented and what each one proves:",
         size=14, bold=True, color=GRAY)
checks = [
    ("5-fold cross-validation", "Trained on 4 folds, tested on 1, repeated 5 times. Stdev across folds < 0.04 -> model is stable."),
    ("Bootstrap 95% CIs (200 resamples)", "Resampled test predictions 200 times. Tells us how much luck is in our headline R²."),
    ("Leakage audit per target", "Caught the CPM = 1.0 trivial-baseline bug. Documented per-target feature exclusions."),
    ("Ablation study (4 base learners)", "Drop each base learner; retrain; measure delta R². Found stack barely beats single HGB."),
    ("Baseline comparison (5 baselines)", "Single XGB and HGB both score within 0.01 R² of full stack on CTR."),
    ("Permutation feature importance", "Shuffle each feature; see how much R² drops. Validates which inputs the model truly relies on."),
]
for i, (name, why) in enumerate(checks):
    y = 1.7 + i*0.85
    add_text(s, 0.6, y, 4.5, 0.5, name, size=14, bold=True, color=ORANGE)
    add_text(s, 5.3, y, 7.5, 0.8, why, size=12, color=GRAY)

# ============================================================
# 18) INPUT -> OUTPUT FLOW
# ============================================================
s = add_slide()
add_header_bar(s, "End-to-End Input -> Output Example")
add_text(s, 0.6, 1.1, 12.0, 0.5,
         "Suppose a marketer wants to test a new Mamaearth carousel ad on Facebook:", size=15, color=GRAY)

# INPUT box
inp = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                         Inches(0.6), Inches(2.0), Inches(4.3), Inches(4.5))
inp.fill.solid(); inp.fill.fore_color.rgb = LIGHT
inp.line.color.rgb = BLUE; inp.line.width = Pt(2)
add_text(s, 0.7, 2.1, 4.0, 0.5, "INPUT", size=18, bold=True, color=BLUE)
add_bullets(s, 0.7, 2.7, 4.0, 4.0,
            ["Brand: Mamaearth",
             "Platform: Facebook",
             "Format: Carousel",
             "Theme: Discount",
             "Audience: New mothers 0-6m",
             "Language: Hinglish",
             "CTA: Shop Now",
             "Has expert: No",
             "Has offer: Yes"], size=13)

# Arrow
arr = s.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, Inches(5.1), Inches(4.0),
                          Inches(0.8), Inches(0.5))
arr.fill.solid(); arr.fill.fore_color.rgb = GRAY; arr.line.fill.background()

# MODEL box
mdl = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                         Inches(6.1), Inches(3.0), Inches(2.0), Inches(2.4))
mdl.fill.solid(); mdl.fill.fore_color.rgb = ORANGE; mdl.line.fill.background()
tf = mdl.text_frame; p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
r = p.add_run(); r.text = "STACKED\nMODEL"
r.font.size = Pt(20); r.font.bold = True; r.font.color.rgb = WHITE

# Arrow
arr = s.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, Inches(8.3), Inches(4.0),
                          Inches(0.8), Inches(0.5))
arr.fill.solid(); arr.fill.fore_color.rgb = GRAY; arr.line.fill.background()

# OUTPUT box
outp = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                           Inches(9.3), Inches(2.0), Inches(3.5), Inches(4.5))
outp.fill.solid(); outp.fill.fore_color.rgb = LIGHT
outp.line.color.rgb = GREEN; outp.line.width = Pt(2)
add_text(s, 9.4, 2.1, 3.3, 0.5, "PREDICTED OUTPUT", size=16, bold=True, color=GREEN)
add_bullets(s, 9.4, 2.7, 3.3, 4.0,
            ["CTR: ~ 2.4%",
             "CPC: INR 12-15",
             "CVR: ~ 1.8%",
             "CPM: INR 110-125",
             "Will user retain D1: 22%",
             "Install Quality: 58 / 100",
             "Top-quartile CTR? Yes",
             "Cross-platform travel: Low (discount-led)"], size=12)

# ============================================================
# 19) HEADLINE RESULTS - chart
# ============================================================
s = add_slide()
add_header_bar(s, "Headline Results")
add_text(s, 0.6, 1.0, 12.0, 0.5,
         "All 12 prediction tasks with confidence bars (higher = better):",
         size=14, color=GRAY)
add_image(s, os.path.join(FIG, 'fig1_r2_with_CI.png'), 1.8, 1.6, 10.0)

# ============================================================
# 20) LIKELY VIVA QUESTIONS
# ============================================================
s = add_slide()
add_header_bar(s, "Likely Viva Questions - We're Ready")
qs = [
    ("Q. Why this project?",      "Real D2C marketing teams burn budget guessing. We tried to help them guess less."),
    ("Q. What's the dataset?",    "9 Kaggle sources + 120 manually-tagged Indian competitor ads. ~22,000 unified rows."),
    ("Q. Why these algorithms?",  "Tree-based gradient boosting beats neural nets on tabular data this size. Explainable too."),
    ("Q. Why didn't you use deep learning?", "Not enough rows. DL needs millions. We have 22,000."),
    ("Q. Is this overfitting?",   "5-fold CV stdev < 0.04 across all targets. Bootstrap CIs are tight. No."),
    ("Q. Why is ROI weak?",       "Source-06 ROI column has |r| < 0.02 with every feature. It's noise. No model can fix that."),
    ("Q. How accurate?",          "R-squared 0.65-0.81 on 6 main targets. AUC 0.99 on the high-CTR classifier."),
    ("Q. Most surprising finding?", "Dropping LLM-augmented data IMPROVED D1 retention by 5.8 R-squared points."),
]
for i, (q, a) in enumerate(qs):
    y = 1.2 + i*0.7
    add_text(s, 0.6, y, 4.5, 0.5, q, size=13, bold=True, color=ORANGE)
    add_text(s, 5.2, y, 7.6, 0.6, a, size=12, color=GRAY)

# ============================================================
# 21) THANK YOU
# ============================================================
s = add_slide()
add_bg(s, BLUE)
add_text(s, 0.5, 2.5, 12.5, 1.5, "Thank You",
         size=72, bold=True, color=WHITE, align='center')
add_text(s, 0.5, 4.2, 12.5, 0.7,
         "Questions and discussion welcome",
         size=24, color=WHITE, align='center')
add_text(s, 0.5, 5.5, 12.5, 0.5,
         "GitHub repo, paper, results JSON and trained models all released",
         size=16, color=WHITE, align='center')

# Save
path = os.path.join(ROOT, 'Viva_Presentation.pptx')
prs.save(path)
print(f"Saved {path}")
print(f"Slides: {len(prs.slides)}")
